from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, Under World!"
subtitle.text = "演習用"

prs.save("test2.pptx")


def count_characters_in_pptx(file_name):
    prs = Presentation(file_name)
    total_characters = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                total_characters += len(shape.text)
    print("Total characters in the presentation:", total_characters)


count_characters_in_pptx("test2.pptx")
